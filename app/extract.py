# SPDX-License-Identifier: MPL-2.0
"""Text extraction for document formats.

A folder of real work files is mostly PDFs, Word documents, slide decks and
spreadsheets, none of which are readable as plain text. Each extractor here
turns one of those into plain text that the chunker and embedder can handle.

Every extractor degrades gracefully: if its library is not installed the file is
skipped with a single explanatory log line rather than failing the ingest, so a
partial install still produces a working index over whatever it can read.
"""

from __future__ import annotations

import csv
import io
import logging
import os
import zipfile
from html.parser import HTMLParser

from .config import CONFIG

log = logging.getLogger("rag.extract")

# Extensions handled here rather than read as plain text.
DOCUMENT_EXTS = {
    ".pdf",
    ".docx", ".pptx", ".xlsx", ".xlsm", ".xls",
    ".odt", ".ods", ".odp",
    ".rtf", ".eml", ".msg",
    ".csv", ".tsv", ".html", ".htm", ".xml",
    # A comic archive: a zip of page images and nothing else. Readable only
    # with OCR on, and skipped with an explanation when it is off.
    ".cbz",
}

# Page images inside a .cbz, in the order a reader would meet them.
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tif", ".tiff"}

# Legacy binary Office formats (.doc, .ppt) have no dependable pure-Python
# reader. They are reported once rather than silently skipped, so nobody
# concludes the index is complete when it is not.
UNSUPPORTED_EXTS = {".doc", ".ppt"}

# One warning per missing library, not one per file.
_warned: set[str] = set()

# Guard against a single spreadsheet or PDF producing a novel's worth of text.
MAX_EXTRACTED_CHARS = 2_000_000


def _missing(library: str, install: str, path: str) -> None:
    if library not in _warned:
        _warned.add(library)
        log.warning(
            "skipping %s files: %s is not installed (pip install %s)",
            os.path.splitext(path)[1],
            library,
            install,
        )


def _cap(text: str) -> str:
    if len(text) > MAX_EXTRACTED_CHARS:
        return text[:MAX_EXTRACTED_CHARS] + "\n[... truncated ...]"
    return text


# Counted rather than logged per file: a folder of scans would otherwise
# produce one warning per document and bury everything else.
_scanned_pages = 0
_scanned_files = 0


def _note_skipped_scan(path: str, scanned: int, total: int) -> None:
    """Record that a document's pages were images and OCR was off.

    Worth a line, because the alternative is an index that silently omits
    whole documents - and a search that finds nothing looks identical whether
    the corpus lacks the answer or the reader could not see it.
    """
    global _scanned_pages, _scanned_files
    _scanned_pages += scanned
    _scanned_files += 1
    if _scanned_files == 1:
        log.warning(
            "%s has %d/%d page(s) with no text layer - they are images. "
            "Set RAG_OCR=1 in rag/.env to read them (slow: seconds per page).",
            path, scanned, total,
        )
    elif _scanned_files % 50 == 0:
        log.warning(
            "%d files / %d pages skipped so far for having no text layer "
            "(RAG_OCR=1 to read them)",
            _scanned_files, _scanned_pages,
        )


def scan_report() -> tuple[int, int]:
    """Files and pages skipped for having no text layer, since start."""
    return _scanned_files, _scanned_pages


def from_cbz(path: str) -> str | None:
    """A comic archive: page images in a zip, readable only through OCR."""
    if not CONFIG.ocr:
        _note_skipped_scan(path, 1, 1)
        return None
    try:
        from PIL import Image
    except ImportError:
        _missing("Pillow", "Pillow", path)
        return None
    from . import ocr as ocr_module

    try:
        with zipfile.ZipFile(path) as archive:
            names = sorted(
                name for name in archive.namelist()
                if os.path.splitext(name)[1].lower() in IMAGE_EXTS
                and not name.endswith("/")
            )
            if not names:
                return None
            if len(names) > CONFIG.archive_max_members:
                names = names[: CONFIG.archive_max_members]
            log.info("OCR: %s - reading %d page image(s)", path, len(names))
            parts: list[str] = []
            for number, name in enumerate(names, start=1):
                try:
                    with archive.open(name) as handle:
                        image = Image.open(io.BytesIO(handle.read())).convert("RGB")
                except Exception as exc:
                    log.debug("could not read %s!%s: %s", path, name, exc)
                    continue
                text = ocr_module.image_text(
                    image,
                    band=CONFIG.ocr_band,
                    overlap=CONFIG.ocr_overlap,
                    min_confidence=CONFIG.ocr_min_confidence,
                )
                if text.strip():
                    parts.append(f"[page {number}]\n{text.strip()}")
    except (zipfile.BadZipFile, OSError) as exc:
        log.warning("could not open %s: %s", path, exc)
        return None
    return _cap("\n\n".join(parts)) if parts else None


def _strip_boilerplate(pages: list[list[str]]) -> list[list[str]]:
    """Remove running headers, footers and page numbers from a PDF.

    A book's every page carries the same footer - "BUILDING AI THAT KNOWS YOUR
    DATA - VOLUME I" and a page number. Extracted verbatim, that text lands in
    every chunk of the document, so a query mentioning any word in it matches
    the whole book uniformly and the actually-relevant page has nothing to
    distinguish it. Stripping repeated lines removes noise the reader never
    sees on the page as content anyway.

    Only lines repeating across most pages are removed, and only in documents
    long enough for "most pages" to mean something.
    """
    if len(pages) < 4:
        return pages

    counts: dict[str, int] = {}
    for lines in pages:
        # Headers and footers live at the edges; a line repeated in the middle
        # of many pages is more likely to be real content.
        #
        # Counted once per PAGE, not once per occurrence: on a page of three
        # lines or fewer the head and tail slices overlap, and counting both
        # would double every line, making a line on two pages look like one on
        # four. The threshold means "appears on N pages".
        edges = {
            line.strip()
            for line in lines[:3] + lines[-3:]
            if line.strip() and len(line.strip()) <= 120
        }
        for text in edges:
            counts[text] = counts.get(text, 0) + 1

    threshold = max(3, int(len(pages) * 0.5))
    boilerplate = {text for text, n in counts.items() if n >= threshold}
    if not boilerplate:
        return pages

    log.debug("stripping %d repeated header/footer line(s)", len(boilerplate))
    cleaned: list[list[str]] = []
    for lines in pages:
        kept = [
            line
            for line in lines
            # A bare page number differs on every page, so it never repeats -
            # match it by shape instead.
            if line.strip() not in boilerplate and not line.strip().isdigit()
        ]
        cleaned.append(kept)
    return cleaned


def from_pdf(path: str) -> str | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        _missing("pypdf", "pypdf", path)
        return None
    try:
        reader = PdfReader(path)
    except Exception as exc:
        log.warning("could not open PDF %s: %s", path, exc)
        return None

    raw_pages: list[list[str]] = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        raw_pages.append(text.split("\n"))

    # Pages whose text layer is empty or near-empty are scanned images. Their
    # words are on the page, just as pixels, and left here they would be
    # indexed as nothing - which reads exactly like a document that does not
    # mention the thing you searched for.
    scanned = {
        number
        for number, lines in enumerate(raw_pages, start=1)
        if len("".join(lines).strip()) < CONFIG.ocr_min_chars
    }
    if scanned:
        if CONFIG.ocr:
            from . import ocr as ocr_module

            log.info("OCR: %s - reading %d page(s) of images", path, len(scanned))
            for number, text in ocr_module.pdf_page_text(path, scanned, CONFIG).items():
                raw_pages[number - 1] = text.split("\n")
        else:
            _note_skipped_scan(path, len(scanned), len(raw_pages))

    pages = raw_pages if os.environ.get("RAG_PDF_KEEP_BOILERPLATE") else _strip_boilerplate(raw_pages)

    parts: list[str] = []
    for number, lines in enumerate(pages, start=1):
        body = "\n".join(lines).strip()
        if body:
            # The page marker survives into the chunk, so a hit can cite a page.
            parts.append(f"[page {number}]\n{body}")
    return _cap("\n\n".join(parts)) if parts else None


def from_docx(path: str) -> str | None:
    try:
        import docx
    except ImportError:
        _missing("python-docx", "python-docx", path)
        return None
    try:
        document = docx.Document(path)
    except Exception as exc:
        log.warning("could not open %s: %s", path, exc)
        return None

    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _cap("\n\n".join(parts)) if parts else None


def from_pptx(path: str) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        _missing("python-pptx", "python-pptx", path)
        return None
    try:
        deck = Presentation(path)
    except Exception as exc:
        log.warning("could not open %s: %s", path, exc)
        return None

    parts: list[str] = []
    for number, slide in enumerate(deck.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
        if lines:
            parts.append(f"[slide {number}]\n" + "\n".join(lines))
    return _cap("\n\n".join(parts)) if parts else None


def from_xlsx(path: str) -> str | None:
    try:
        import openpyxl
    except ImportError:
        _missing("openpyxl", "openpyxl", path)
        return None
    try:
        book = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:
        log.warning("could not open %s: %s", path, exc)
        return None

    parts: list[str] = []
    try:
        for sheet in book.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[sheet {sheet.title}]\n" + "\n".join(rows))
    finally:
        book.close()
    return _cap("\n\n".join(parts)) if parts else None


def from_delimited(path: str) -> str | None:
    delimiter = "\t" if path.lower().endswith(".tsv") else ","
    try:
        with open(path, newline="", encoding="utf-8", errors="replace") as handle:
            rows = [
                " | ".join(c.strip() for c in row if c.strip())
                for row in csv.reader(handle, delimiter=delimiter)
            ]
    except OSError as exc:
        log.warning("could not read %s: %s", path, exc)
        return None
    body = "\n".join(r for r in rows if r)
    return _cap(body) if body.strip() else None


class _TextHTML(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._skip:
            self._skip -= 1

    def handle_data(self, data):
        if not self._skip and data.strip():
            self.parts.append(data.strip())


def from_html(path: str) -> str | None:
    try:
        with open(path, encoding="utf-8", errors="replace") as handle:
            raw = handle.read()
    except OSError as exc:
        log.warning("could not read %s: %s", path, exc)
        return None
    parser = _TextHTML()
    try:
        parser.feed(raw)
    except Exception:
        return None
    body = "\n".join(parser.parts)
    return _cap(body) if body.strip() else None


def from_xls(path: str) -> str | None:
    """Legacy .xls. xlrd 2.x dropped .xlsx but still reads the old format."""
    try:
        import xlrd
    except ImportError:
        _missing("xlrd", "xlrd", path)
        return None
    try:
        book = xlrd.open_workbook(path)
    except Exception as exc:
        log.warning("could not open %s: %s", path, exc)
        return None

    parts: list[str] = []
    for sheet in book.sheets():
        rows: list[str] = []
        for index in range(sheet.nrows):
            cells = [str(c.value).strip() for c in sheet.row(index) if str(c.value).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[sheet {sheet.name}]\n" + "\n".join(rows))
    return _cap("\n\n".join(parts)) if parts else None


def from_opendocument(path: str) -> str | None:
    """ODF (.odt/.ods/.odp) is a zip with an XML payload - stdlib is enough."""
    import xml.etree.ElementTree as ET
    import zipfile

    try:
        with zipfile.ZipFile(path) as archive:
            raw = archive.read("content.xml")
    except (OSError, KeyError, zipfile.BadZipFile) as exc:
        log.warning("could not open %s: %s", path, exc)
        return None
    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        log.warning("could not parse %s: %s", path, exc)
        return None

    parts = [t.strip() for t in root.itertext() if t and t.strip()]
    body = "\n".join(parts)
    return _cap(body) if body.strip() else None


def from_rtf(path: str) -> str | None:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        _missing("striprtf", "striprtf", path)
        return None
    try:
        with open(path, encoding="utf-8", errors="replace") as handle:
            body = rtf_to_text(handle.read(), errors="ignore")
    except Exception as exc:
        log.warning("could not read %s: %s", path, exc)
        return None
    return _cap(body) if body and body.strip() else None


def from_eml(path: str) -> str | None:
    """Saved email. Headers matter for retrieval, so they are kept."""
    import email
    from email import policy

    try:
        with open(path, "rb") as handle:
            message = email.message_from_binary_file(handle, policy=policy.default)
    except Exception as exc:
        log.warning("could not read %s: %s", path, exc)
        return None

    header_lines = [
        f"{name}: {message.get(name)}"
        for name in ("From", "To", "Cc", "Date", "Subject")
        if message.get(name)
    ]
    try:
        part = message.get_body(preferencelist=("plain", "html"))
        body = part.get_content() if part else ""
    except Exception:
        body = ""
    if body and part is not None and part.get_content_type() == "text/html":
        parser = _TextHTML()
        parser.feed(body)
        body = "\n".join(parser.parts)

    text = "\n".join(header_lines) + "\n\n" + (body or "")
    return _cap(text) if text.strip() else None


def from_msg(path: str) -> str | None:
    """Outlook .msg - very common in shared drives."""
    try:
        import extract_msg
    except ImportError:
        _missing("extract-msg", "extract-msg", path)
        return None
    try:
        message = extract_msg.Message(path)
    except Exception as exc:
        log.warning("could not open %s: %s", path, exc)
        return None
    try:
        header_lines = [
            f"{label}: {value}"
            for label, value in (
                ("From", message.sender),
                ("To", message.to),
                ("Cc", message.cc),
                ("Date", message.date),
                ("Subject", message.subject),
            )
            if value
        ]
        text = "\n".join(header_lines) + "\n\n" + (message.body or "")
    finally:
        try:
            message.close()
        except Exception:
            pass
    return _cap(text) if text.strip() else None


def from_xml(path: str) -> str | None:
    import xml.etree.ElementTree as ET

    try:
        root = ET.parse(path).getroot()
    except (OSError, ET.ParseError) as exc:
        log.warning("could not parse %s: %s", path, exc)
        return None
    body = "\n".join(t.strip() for t in root.itertext() if t and t.strip())
    return _cap(body) if body.strip() else None


_EXTRACTORS = {
    ".pdf": from_pdf,
    ".xls": from_xls,
    ".odt": from_opendocument,
    ".ods": from_opendocument,
    ".odp": from_opendocument,
    ".rtf": from_rtf,
    ".eml": from_eml,
    ".msg": from_msg,
    ".xml": from_xml,
    ".docx": from_docx,
    ".pptx": from_pptx,
    ".xlsx": from_xlsx,
    ".xlsm": from_xlsx,
    ".csv": from_delimited,
    ".tsv": from_delimited,
    ".html": from_html,
    ".htm": from_html,
    ".cbz": from_cbz,
}


def is_document(path: str) -> bool:
    return os.path.splitext(path)[1].lower() in DOCUMENT_EXTS


def extract(path: str) -> str | None:
    """Extract plain text from a document, or None if it cannot be read."""
    extractor = _EXTRACTORS.get(os.path.splitext(path)[1].lower())
    if extractor is None:
        return None
    try:
        return extractor(path)
    except Exception as exc:  # a corrupt file must not abort the whole ingest
        log.warning("extraction failed for %s: %s", path, exc)
        return None
